import os
import fitz  # PyMuPDF
import pptx
import random

def extract_text_from_file(file_path):
    ext = file_path.split('.')[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext in ["ppt", "pptx"]:
        return extract_text_from_pptx(file_path)
    return ""

def extract_text_from_pdf(file_path):
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        text = f"Error reading PDF: {e}"
    return text

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = pptx.Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"Error reading PPTX: {e}"
    return text

def generate_mcq(text, num_questions=3):
    if not text.strip():
        return []
    questions = []
    sentences = [s.strip() for s in text.split(".") if len(s.strip().split()) > 5]
    random.shuffle(sentences)
    for i in range(min(num_questions, len(sentences))):
        q = sentences[i]
        questions.append({
            "question": f"What is the implication of: {q}?",
            "options": ["Option A", "Option B", "Option C", "Option D"],
            "correct": 0
        })
    return questions